import json
import zipfile
from io import BytesIO
from pathlib import Path
import html
from docx import Document
from docx.shared import Inches, Pt, RGBColor
from docx.enum.text import WD_ALIGN_PARAGRAPH

STATIC_DIR = Path(__file__).resolve().parents[1] / "static" / "videos"


def lesson_to_markdown(data: dict) -> str:
    lines = [f"# {data.get('topic', 'Lesson')}", ""]
    if data.get("overview"):
        lines += ["## Overview", data["overview"], ""]
    if data.get("learning_objectives"):
        lines += ["## Learning Objectives", *[f"- {item}" for item in data["learning_objectives"]], ""]
    for section in data.get("sections") or []:
        lines += [f"## {section.get('heading', 'Section')}", section.get("content", "")]
        if section.get("example"):
            lines += [f"*Example:* {section['example']}"]
        lines.append("")
    if data.get("summary"):
        lines += ["## Summary", data["summary"], ""]
    return "\n".join(lines)


def flashcards_to_markdown(data: dict) -> str:
    lines = [f"# {data.get('title', 'Flashcards')}", ""]
    for index, card in enumerate(data.get("cards") or [], start=1):
        lines += [f"## Card {index}", f"**Front:** {card.get('front', '')}", f"**Back:** {card.get('back', '')}", ""]
    return "\n".join(lines)


def quiz_to_markdown(data: dict) -> str:
    lines = [f"# {data.get('title', 'Quiz')}", f"Difficulty: {data.get('difficulty', 'Medium')}", ""]
    for index, question in enumerate(data.get("questions") or [], start=1):
        lines += [f"## Question {index}", question.get("question", "")]
        for option in question.get("options") or []:
            lines.append(f"- {option}")
        lines += [f"**Answer:** {question.get('answer', '')}", ""]
    return "\n".join(lines)


def build_word_document_bytes(title: str, sections: list[tuple[str, str | list[str]]]) -> BytesIO:
    body = [
        "<html><head><meta charset='utf-8'>",
        "<style>",
        "body{font-family:Calibri,Arial,sans-serif;padding:28px;color:#1e293b;}",
        "h1{color:#0f766e;margin-bottom:6px;}",
        "h2{color:#155e75;margin-top:22px;}",
        "p,li{line-height:1.7;font-size:12pt;}",
        ".card{margin:14px 0;padding:14px 16px;border:1px solid #cbd5e1;border-radius:14px;background:#f8fafc;}",
        "</style></head><body>",
        f"<h1>{html.escape(title)}</h1>",
    ]
    for heading, content in sections:
        body.append(f"<h2>{html.escape(heading)}</h2>")
        if isinstance(content, list):
            body.append("<ul>")
            body.extend(f"<li>{html.escape(str(item))}</li>" for item in content if item)
            body.append("</ul>")
        else:
            paragraphs = [segment.strip() for segment in str(content).split("\n") if segment.strip()]
            if not paragraphs:
                body.append("<p></p>")
            for paragraph in paragraphs:
                body.append(f"<p>{html.escape(paragraph)}</p>")
    body.append("</body></html>")
    buffer = BytesIO("".join(body).encode("utf-8"))
    buffer.seek(0)
    return buffer


def build_video_zip(project_id: int) -> BytesIO | None:
    project_dir = STATIC_DIR / f"video-{project_id}"
    if not project_dir.exists():
        return None
    buffer = BytesIO()
    with zipfile.ZipFile(buffer, "w", zipfile.ZIP_DEFLATED) as archive:
        for file_path in project_dir.rglob("*"):
            if file_path.is_file():
                archive.write(file_path, arcname=file_path.relative_to(project_dir))
    buffer.seek(0)
    return buffer


def build_json_bytes(data: dict) -> BytesIO:
    buffer = BytesIO(json.dumps(data, ensure_ascii=False, indent=2).encode("utf-8"))
    buffer.seek(0)
    return buffer


def build_true_docx_bytes(title: str, content_dict: dict, content_type: str = "lesson") -> BytesIO:
    """Create a native .docx document using python-docx"""
    doc = Document()

    # Set narrower margins for more content per page
    for section in doc.sections:
        section.top_margin = Inches(1)
        section.bottom_margin = Inches(1)
        section.left_margin = Inches(1.2)
        section.right_margin = Inches(1.2)

    def add_colored_heading(text: str, level: int = 1, color: RGBColor = RGBColor(15, 118, 110)):
        h = doc.add_heading(text, level=level)
        h.paragraph_format.space_before = Pt(16)
        h.paragraph_format.space_after = Pt(6)
        for run in h.runs:
            run.font.color.rgb = color
        return h

    def add_body(text: str):
        p = doc.add_paragraph(text)
        p.paragraph_format.space_after = Pt(8)
        p.paragraph_format.line_spacing = Pt(14)
        for run in p.runs:
            run.font.size = Pt(11)
        return p

    if content_type == "lesson":
        # ── Cover block ──
        title_para = doc.add_heading(title, 0)
        title_para.alignment = WD_ALIGN_PARAGRAPH.CENTER
        title_para.paragraph_format.space_after = Pt(4)
        for run in title_para.runs:
            run.font.size = Pt(26)
            run.font.color.rgb = RGBColor(15, 118, 110)
            run.font.bold = True

        # Subtitle / metadata
        meta_lines = []
        if content_dict.get("language"):
            meta_lines.append(f"Language: {content_dict['language']}")
        style_val = content_dict.get("style") or content_dict.get("teaching_style")
        if style_val:
            meta_lines.append(f"Style: {style_val}")
        if content_dict.get("duration"):
            meta_lines.append(f"Duration: {content_dict['duration']} minutes")
        if meta_lines:
            meta_p = doc.add_paragraph(" · ".join(meta_lines))
            meta_p.alignment = WD_ALIGN_PARAGRAPH.CENTER
            meta_p.paragraph_format.space_after = Pt(18)
            for run in meta_p.runs:
                run.font.size = Pt(10)
                run.font.color.rgb = RGBColor(100, 116, 139)
                run.font.italic = True

        doc.add_paragraph()  # spacer

        # Overview
        if content_dict.get("overview"):
            add_colored_heading("Overview")
            add_body(content_dict["overview"])

        # Learning Objectives
        if content_dict.get("learning_objectives"):
            add_colored_heading("Learning Objectives")
            for obj in content_dict["learning_objectives"]:
                p = doc.add_paragraph(style="List Bullet")
                run = p.add_run(obj)
                run.font.size = Pt(11)
                p.paragraph_format.space_after = Pt(4)

            doc.add_paragraph()  # spacer after list

        # Content sections
        for section in content_dict.get("sections") or []:
            add_colored_heading(section.get("heading", "Section"), level=1, color=RGBColor(21, 94, 117))

            content_text = section.get("content", "")
            if content_text:
                # Split into paragraphs for better formatting
                for para_text in content_text.split("\n"):
                    if para_text.strip():
                        add_body(para_text.strip())

            if section.get("example"):
                example_p = doc.add_paragraph()
                example_p.paragraph_format.left_indent = Inches(0.4)
                example_p.paragraph_format.space_before = Pt(4)
                example_p.paragraph_format.space_after = Pt(12)
                label_run = example_p.add_run("💡 Example: ")
                label_run.font.bold = True
                label_run.font.color.rgb = RGBColor(15, 118, 110)
                label_run.font.size = Pt(11)
                content_run = example_p.add_run(section["example"])
                content_run.font.italic = True
                content_run.font.size = Pt(11)

        # Key Points
        if content_dict.get("key_points"):
            add_colored_heading("Key Points", level=1, color=RGBColor(14, 116, 144))
            for kp in content_dict["key_points"]:
                p = doc.add_paragraph(style="List Bullet")
                run = p.add_run(kp)
                run.font.size = Pt(11)
                p.paragraph_format.space_after = Pt(4)
            doc.add_paragraph()

        # Slides summary
        if content_dict.get("slides"):
            add_colored_heading("Slide Outlines", level=1, color=RGBColor(14, 116, 144))
            for i, slide in enumerate(content_dict["slides"], start=1):
                slide_h = doc.add_heading(f"Slide {i}: {slide.get('title', '')}", level=2)
                slide_h.paragraph_format.space_before = Pt(10)
                slide_h.paragraph_format.space_after = Pt(3)
                for run in slide_h.runs:
                    run.font.color.rgb = RGBColor(51, 65, 85)
                if slide.get("content"):
                    add_body(slide["content"])

        # Summary
        if content_dict.get("summary"):
            add_colored_heading("Summary")
            add_body(content_dict["summary"])

    elif content_type == "quiz":
        title_para = doc.add_heading(title, 0)
        title_para.alignment = WD_ALIGN_PARAGRAPH.CENTER
        for run in title_para.runs:
            run.font.size = Pt(24)
            run.font.color.rgb = RGBColor(15, 118, 110)

        diff_p = doc.add_paragraph(f"Difficulty: {content_dict.get('difficulty', 'Medium')}")
        diff_p.paragraph_format.space_after = Pt(16)
        for run in diff_p.runs:
            run.font.bold = True
            run.font.size = Pt(12)

        for index, question in enumerate(content_dict.get("questions") or [], start=1):
            h3 = doc.add_heading(f"Question {index}", level=2)
            h3.paragraph_format.space_before = Pt(14)
            for run in h3.runs:
                run.font.color.rgb = RGBColor(30, 64, 175)

            q_p = doc.add_paragraph(question.get("question", ""))
            q_p.paragraph_format.space_after = Pt(6)

            for option in question.get("options") or []:
                opt_p = doc.add_paragraph(style="List Bullet")
                opt_p.add_run(option).font.size = Pt(11)

            answer_p = doc.add_paragraph()
            answer_p.paragraph_format.space_before = Pt(4)
            a_label = answer_p.add_run("Answer: ")
            a_label.font.bold = True
            a_label.font.color.rgb = RGBColor(22, 163, 74)
            a_val = answer_p.add_run(question.get("answer", ""))
            a_val.font.bold = True
            a_val.font.color.rgb = RGBColor(22, 163, 74)

            if question.get("explanation"):
                exp_p = doc.add_paragraph(f"Explanation: {question['explanation']}")
                exp_p.paragraph_format.space_after = Pt(8)
                for run in exp_p.runs:
                    run.font.italic = True
                    run.font.color.rgb = RGBColor(100, 116, 139)

    elif content_type == "flashcards":
        title_para = doc.add_heading(title, 0)
        title_para.alignment = WD_ALIGN_PARAGRAPH.CENTER
        for run in title_para.runs:
            run.font.color.rgb = RGBColor(15, 118, 110)

        for index, card in enumerate(content_dict.get("cards") or [], start=1):
            h3 = doc.add_heading(f"Card {index}", level=2)
            h3.paragraph_format.space_before = Pt(14)
            for run in h3.runs:
                run.font.color.rgb = RGBColor(30, 64, 175)

            front_p = doc.add_paragraph()
            fl = front_p.add_run("Front: ")
            fl.font.bold = True
            fl.font.color.rgb = RGBColor(15, 118, 110)
            front_p.add_run(card.get("front", "")).font.size = Pt(11)

            back_p = doc.add_paragraph()
            bl = back_p.add_run("Back: ")
            bl.font.bold = True
            bl.font.color.rgb = RGBColor(21, 94, 117)
            back_p.add_run(card.get("back", "")).font.size = Pt(11)
            back_p.paragraph_format.space_after = Pt(12)


    
    buffer = BytesIO()
    doc.save(buffer)
    buffer.seek(0)
    return buffer


def build_pptx_bytes(title: str, content_dict: dict) -> BytesIO:
    """Generate a PowerPoint presentation using python-pptx."""
    from pptx import Presentation
    from pptx.util import Inches, Pt, Emu
    from pptx.dml.color import RGBColor as PptxRGB
    from pptx.enum.text import PP_ALIGN

    TEAL = PptxRGB(0x0F, 0x76, 0x6E)
    DARK = PptxRGB(0x1E, 0x29, 0x3B)
    BODY_CLR = PptxRGB(0x33, 0x41, 0x55)
    WHITE = PptxRGB(0xFF, 0xFF, 0xFF)
    LIGHT_BG = PptxRGB(0xF1, 0xF5, 0xF9)

    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)

    blank_layout = prs.slide_layouts[6]  # Blank

    def add_slide(slide_title_text: str, body_lines: list[str], is_title_slide: bool = False):
        slide = prs.slides.add_slide(blank_layout)

        # Background
        bg = slide.background
        fill = bg.fill
        fill.solid()
        fill.fore_color.rgb = WHITE if not is_title_slide else DARK

        title_top = Inches(0.4)
        title_height = Inches(1.2)
        body_top = Inches(1.8)
        body_height = Inches(5.3)
        left = Inches(0.6)
        width = Inches(12.1)

        # Title box
        txBox = slide.shapes.add_textbox(left, title_top, width, title_height)
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.LEFT
        run = p.add_run()
        run.text = slide_title_text
        run.font.size = Pt(36) if is_title_slide else Pt(28)
        run.font.bold = True
        run.font.color.rgb = WHITE if is_title_slide else TEAL

        # Divider line (thin rectangle)
        left_emu = int(left)
        line_top = int(title_top + title_height - Inches(0.05))
        line = slide.shapes.add_shape(
            1,  # MSO_SHAPE_TYPE.RECTANGLE
            left_emu, line_top,
            int(width), int(Inches(0.04))
        )
        line.fill.solid()
        line.fill.fore_color.rgb = TEAL
        line.line.fill.background()

        # Body text box
        body_box = slide.shapes.add_textbox(left, body_top, width, body_height)
        tf2 = body_box.text_frame
        tf2.word_wrap = True
        first = True
        for line_text in body_lines:
            if first:
                p2 = tf2.paragraphs[0]
                first = False
            else:
                p2 = tf2.add_paragraph()
            p2.space_after = Pt(6)
            run2 = p2.add_run()
            run2.text = line_text
            run2.font.size = Pt(18)
            run2.font.color.rgb = WHITE if is_title_slide else BODY_CLR

    sections = content_dict.get("sections") or []

    # Slide 1: Title
    subtitle_lines = [
        f"Subject: {content_dict.get('topic', title)}",
        f"Style: {content_dict.get('style') or content_dict.get('teaching_style', 'Interactive')}",
        f"Duration: {content_dict.get('duration', 5)} minutes",
    ]
    add_slide(title, subtitle_lines, is_title_slide=True)

    # Slide 2: Overview
    if content_dict.get("overview"):
        add_slide("Overview", [content_dict["overview"]])

    # Slide 3: Learning Objectives
    if content_dict.get("learning_objectives"):
        objs = [f"• {obj}" for obj in content_dict["learning_objectives"]]
        add_slide("Learning Objectives", objs)

    # Content sections
    for sec in sections:
        lines = []
        if sec.get("content"):
            lines += [l.strip() for l in sec["content"].split("\n") if l.strip()]
        if sec.get("example"):
            lines.append("")
            lines.append(f"Example: {sec['example']}")
        add_slide(sec.get("heading", "Section"), lines)

    # Summary
    if content_dict.get("summary"):
        add_slide("Summary", [content_dict["summary"]])

    # Key points
    if content_dict.get("key_points"):
        kp = [f"• {kp}" for kp in content_dict["key_points"]]
        add_slide("Key Points", kp)

    buffer = BytesIO()
    prs.save(buffer)
    buffer.seek(0)
    return buffer
