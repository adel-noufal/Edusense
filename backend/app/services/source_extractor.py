"""Extract plain text from instructor-uploaded session resources for AI context."""
from pathlib import Path


def _read_text_file(path: Path, limit: int = 12000) -> str:
    for encoding in ("utf-8", "utf-16", "latin-1"):
        try:
            return path.read_text(encoding=encoding)[:limit]
        except Exception:
            continue
    return ""


def _read_docx(path: Path, limit: int = 12000) -> str:
    try:
        from docx import Document
        doc = Document(str(path))
        parts = [p.text.strip() for p in doc.paragraphs if p.text.strip()]
        return "\n".join(parts)[:limit]
    except Exception:
        return ""


def _read_pptx(path: Path, limit: int = 12000) -> str:
    try:
        from pptx import Presentation
        prs = Presentation(str(path))
        parts = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    parts.append(shape.text.strip())
        return "\n".join(parts)[:limit]
    except Exception:
        return ""


def _read_pdf(path: Path, limit: int = 12000) -> str:
    try:
        from pypdf import PdfReader
        reader = PdfReader(str(path))
        parts = []
        for page in reader.pages[:20]:
            text = page.extract_text() or ""
            if text.strip():
                parts.append(text.strip())
        return "\n".join(parts)[:limit]
    except Exception:
        return ""


def extract_text_from_file(file_path: str, file_type: str = "") -> str:
    path = Path(file_path)
    if not path.exists():
        return ""
    suffix = path.suffix.lower()
    mime = (file_type or "").lower()
    if suffix in {".txt", ".md", ".csv", ".json"}:
        return _read_text_file(path)
    if suffix in {".docx", ".doc"} or "word" in mime:
        return _read_docx(path)
    if suffix in {".pptx", ".ppt"} or "presentation" in mime or "powerpoint" in mime:
        return _read_pptx(path)
    if suffix == ".pdf" or "pdf" in mime:
        return _read_pdf(path)
    return ""


def build_session_source_context(resources: list) -> str:
    """Combine extracted text from multiple session resources."""
    chunks = []
    for resource in resources:
        text = extract_text_from_file(resource.file_path, getattr(resource, "file_type", "") or "")
        if text.strip():
            name = getattr(resource, "name", Path(resource.file_path).name)
            chunks.append(f"--- Source: {name} ---\n{text.strip()}")
    if not chunks:
        return ""
    combined = "\n\n".join(chunks)
    return combined[:16000]
